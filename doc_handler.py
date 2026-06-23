#!/usr/bin/env python3
"""
doc_handler — tag academic documents with a local LLM, then move by prefix.

Tiered classification (handles no-text / handwritten / scanned / book PDFs):
  1. TEXT      first 2 pages -> text model.
  2. ESCALATE  if result is 99UNS (e.g. page 1 was a book cover/preface/TOC),
               re-read up to 5 pages and retry before giving up.
  3. VISION    no extractable text -> render page 1 to image -> vision model
               (page 3 retry if still 99UNS).
  4. FILENAME  last resort -> filename + folder path.

Output prefix: `[STREAM-SUBJECT] name.pdf`. Vocabulary + system prompt live in
`system_prompt.md`. Dry-run by default; --apply renames; --move relocates.

  python doc_handler.py "D:\\AcademicsCOPY" --vision --vision-model qwen2-vl-7b   # dry-run
  python doc_handler.py "D:\\AcademicsCOPY" --apply --vision --vision-model qwen2-vl-7b
  python doc_handler.py "D:\\AcademicsCOPY" --move "D:\\Archive\\Academics"        # dry-run move
  python doc_handler.py "D:\\AcademicsCOPY" --move "D:\\Archive\\Academics" --apply

Deps: see requirements.txt  (pip install -r requirements.txt)
"""
from __future__ import annotations
import os, sys, json, csv, re, base64, argparse, urllib.request

HERE=os.path.dirname(os.path.abspath(__file__))
EXT_TEXT={".pdf",".txt",".md",".docx",".pptx",".ppt",".doc"}
STREAMS={"CW","GATE","PROJ","RES","REC","REF"}
TYPES={"notes","pyq","book","slides","assignment","lab","report","datasheet","syllabus","solution","misc"}
VALID={"00MM","01CA","02SEMI","03PN","04BJT","05MOS","06OPAMP","07ANLG","08DIG",
       "09SNS","10CTRL","11COMM","12EMAG","13TOOLS","90HUM","NA","99UNS"}
MIN_TEXT=80
DEEP_PAGES=5
DEEP_CAP=4000

def load_sys(path):
    try:
        s=open(path,encoding="utf-8").read()
        m=re.search(r"<SYSTEM>(.*?)</SYSTEM>",s,re.S)
        if m: return m.group(1).strip()
    except Exception: pass
    return "Reply: STREAM SUBJECT TYPE CONF. Subjects:"+",".join(sorted(VALID))

def pdf_text(path,pages=2,cap=2000):
    try:
        import fitz; d=fitz.open(path); t=""
        for i in range(min(pages,d.page_count)):
            t+=d[i].get_text()
            if len(t)>cap: break
        d.close(); return t[:cap]
    except Exception: return ""

def doc_text(path,cap=2000):
    e=os.path.splitext(path)[1].lower()
    try:
        if e==".pdf": return pdf_text(path,2,cap)
        if e in (".txt",".md"): return open(path,encoding="utf-8",errors="replace").read()[:cap]
        if e==".docx":
            import docx; return "\n".join(p.text for p in docx.Document(path).paragraphs)[:cap]
        if e==".pptx":
            from pptx import Presentation
            return " ".join(s.text for sl in Presentation(path).slides for s in sl.shapes
                             if getattr(s,'has_text_frame',False))[:cap]
    except Exception: return ""
    return ""

def page_png(path,page=0,dpi=120):
    try:
        import fitz; d=fitz.open(path)
        if d.page_count<=page: page=0
        if d.page_count==0: return None
        b=d[page].get_pixmap(dpi=dpi).tobytes("png"); d.close(); return b
    except Exception: return None

def _post(api,payload):
    req=urllib.request.Request(api,data=json.dumps(payload).encode(),
                               headers={"Content-Type":"application/json"})
    r=json.load(urllib.request.urlopen(req,timeout=180))
    return r["choices"][0]["message"]["content"]

def parse(out):
    o=(out or "").upper()
    stream=next((s for s in sorted(STREAMS,key=len,reverse=True) if re.search(r'\b'+s+r'\b',o)),"CW")
    subj=next((c for c in sorted(VALID,key=len,reverse=True) if c in o),"99UNS")
    typ=next((t for t in TYPES if t.upper() in o),"misc")
    conf="high" if "HIGH" in o else "low"
    return stream,subj,typ,conf

def cls_text(api,model,sysp,name,rel,snip):
    usr=f"Filename: {name}\nFolder: {rel}\nText:\n{snip[:DEEP_CAP]}\n\nAnswer (STREAM SUBJECT TYPE CONF):"
    try: return parse(_post(api,{"model":model,"temperature":0,"max_tokens":16,
        "messages":[{"role":"system","content":sysp},{"role":"user","content":usr}]}))
    except Exception: return ("CW","99UNS","misc","low")

def cls_vision(api,model,sysp,name,rel,png,note=""):
    b64="data:image/png;base64,"+base64.b64encode(png).decode()
    usr=[{"type":"text","text":f"Filename: {name}\nFolder: {rel}\n{note}This is a page (may be handwritten/scanned). Answer (STREAM SUBJECT TYPE CONF):"},
         {"type":"image_url","image_url":{"url":b64}}]
    try: return parse(_post(api,{"model":model,"temperature":0,"max_tokens":16,
        "messages":[{"role":"system","content":sysp},{"role":"user","content":usr}]}))
    except Exception: return ("CW","99UNS","misc","low")

def classify(a,sysp,full,fn,rel):
    """Returns (stream,subj,typ,conf,source). Escalates on 99UNS."""
    ispdf=full.lower().endswith(".pdf")
    snip=doc_text(full)
    if len(snip.strip())>=MIN_TEXT:
        st,su,ty,cf=cls_text(a.api,a.model,sysp,fn,rel,snip); src="text"
        if su=="99UNS" and ispdf:                       # ESCALATE: book cover/preface -> read 5 pages
            deep=pdf_text(full,DEEP_PAGES,DEEP_CAP)
            if len(deep)>len(snip):
                st2,su2,ty2,cf2=cls_text(a.api,a.model,sysp,fn,rel,deep)
                if su2!="99UNS": return st2,su2,ty2,cf2,"text5"
        return st,su,ty,cf,src
    if a.vision and ispdf and (png:=page_png(full,0)):
        st,su,ty,cf=cls_vision(a.api,a.vision_model,sysp,fn,rel,png); src="vision"
        if su=="99UNS" and (png3:=page_png(full,2)):    # ESCALATE: try a mid page
            st2,su2,ty2,cf2=cls_vision(a.api,a.vision_model,sysp,fn,rel,png3,"(page 3) ")
            if su2!="99UNS": return st2,su2,ty2,cf2,"vision3"
        return st,su,ty,cf,src
    st,su,ty,cf=cls_text(a.api,a.model,sysp,fn,rel,"")   # filename only
    return st,su,ty,cf,"filename"

def move_by_prefix(root,dest,apply):
    pat=re.compile(r'^\[([A-Z]+)-([0-9A-Z]+)\]\s+(.*)$'); n=0; rows=[]
    for dp,_,fns in os.walk(root):
        for fn in fns:
            m=pat.match(fn)
            if not m: continue
            stream,subj,base=m.group(1),m.group(2),m.group(3)
            newp=os.path.join(dest,stream,subj,base); rows.append((os.path.join(dp,fn),newp)); n+=1
            if apply:
                os.makedirs(os.path.join(dest,stream,subj),exist_ok=True)
                try: os.rename(os.path.join(dp,fn),newp)
                except Exception as e: print("move-fail:",fn,e)
    log=os.path.join(dest if apply else root,"_move_log.csv")
    try:
        os.makedirs(os.path.dirname(log),exist_ok=True)
        with open(log,"w",encoding="utf-8",newline="") as g: csv.writer(g).writerows([("from","to")]+rows)
    except Exception: pass
    print(f"{'MOVED' if apply else 'DRY-RUN move'}: {n} files -> {dest}")

def main():
    ap=argparse.ArgumentParser()
    ap.add_argument("root")
    ap.add_argument("--api",default="http://localhost:1234/v1/chat/completions")
    ap.add_argument("--model",default="local-model")
    ap.add_argument("--vision",action="store_true")
    ap.add_argument("--vision-model",default="local-vision")
    ap.add_argument("--prompt",default=os.path.join(HERE,"system_prompt.md"))
    ap.add_argument("--apply",action="store_true")
    ap.add_argument("--move",default=None)
    ap.add_argument("--log",default=None)
    a=ap.parse_args()
    if a.move: move_by_prefix(a.root,a.move,a.apply); return
    sysp=load_sys(a.prompt); rows=[]; n=0
    for dp,_,fns in os.walk(a.root):
        for fn in fns:
            if os.path.splitext(fn)[1].lower() not in EXT_TEXT or fn.startswith("["): continue
            full=os.path.join(dp,fn); rel=os.path.relpath(dp,a.root)
            st,su,ty,cf,src=classify(a,sysp,full,fn,rel)
            new=f"[{st}-{su}] {fn}"; rows.append((full,fn,new,st,su,ty,cf,src)); n+=1
            print(f"{st:5}{su:7}{ty:10}{cf:5}{src:9}{fn[:48]}")
            if a.apply:
                try: os.rename(full,os.path.join(dp,new))
                except Exception as e: print("  rename-fail:",e)
    logp=a.log or os.path.join(a.root,"_doc_handler_log.csv")
    with open(logp,"w",encoding="utf-8",newline="") as g:
        w=csv.writer(g); w.writerow(["path","old","new","stream","subject","type","conf","source"]); w.writerows(rows)
    print(f"\n{'APPLIED' if a.apply else 'DRY-RUN'}: {n} files. log: {logp}")
    if not a.apply: print("Review log (focus conf=low, source=filename/vision), then --apply.")
if __name__=="__main__": main()
